from openai import OpenAI
from dotenv import load_dotenv
import os
import base64
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
import json


load_dotenv()

#===OpenAI===
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY)
#===Local Files Root directory===
UPLOAD_FOLDER = os.getenv("UPLOAD_ROOT")

#===File-Processing Helper Methods===
DOC_EXTENSIONS = ['.pdf', '.docx', '.pptx', '.txt']
IMG_EXTENSIONS = ['.jpg', '.jpeg', '.png']
def extract_text(file_path, chunk_size=500):
    ext = file_path.lower()
    full_text = ""
    if ext.endswith(".pdf"):
        reader = PdfReader(file_path)
        full_text = "\n".join([page.extract_text() or "" for page in reader.pages])
    elif ext.endswith(".docx"):
        doc = Document(file_path)
        full_text = "\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])
    elif ext.endswith(".pptx"):
        prs = Presentation(file_path)
        full_text = "\n".join([" ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]) for slide in prs.slides])
    elif ext.endswith(".txt"):
        with open(file_path, "r", encoding="utf-8") as f:
            full_text = f.read()
    else:
        raise ValueError("Unsupported file format.")
    return chunk_text(full_text, chunk_size)

def extract_images_from_pdf(document_dir, file_path, images_dir):
    reader = PdfReader(file_path)
    images = []
    alt_text_map = []

    for i, page in enumerate(reader.pages):
        for img_index, image in enumerate(page.images):
            img_data = image.data
            img_filename = f"page-{i}-{img_index}.jpg"
            img_path = os.path.join(images_dir, img_filename)

            with open(img_path, "wb") as f:
                f.write(img_data)

            images.append(img_path)

            alt_text = generate_gpt4_description(img_path)
            alt_text_map.append({
                "path": img_path,
                "alt_text": alt_text
            })

    # Save alt-image map for the PDF
    if alt_text_map:
        map_file_path = os.path.join(document_dir, "alt_image_map.json")
        with open(map_file_path, "w", encoding="utf-8") as f:
            json.dump(alt_text_map, f, indent=4)

    return images


def extract_images_from_docx(document_dir, file_path, images_dir):
    doc = Document(file_path)
    os.makedirs(images_dir, exist_ok=True)
    images = []
    alt_text_map = []

    for shape in doc.inline_shapes:
        alt_text = shape._inline.docPr.get("descr")
        if not alt_text or ':' not in alt_text:
            continue

        image_name, _, alt_content = alt_text.partition(':')
        image_name = image_name.strip()
        alt_content = alt_content.strip()

        if not image_name or not alt_content:
            continue

        try:
            image_name = image_name.replace(" ","_")
            rel_id = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            image_part = doc.part.related_parts[rel_id]
            image_bytes = image_part.blob

            ext = "jpg"
            img_filename = f"{image_name}.{ext}"
            img_path = os.path.join(images_dir, img_filename)

            with open(img_path, "wb") as f:
                f.write(image_bytes)

            images.append(img_path)
            alt_text_map.append({
                "path": img_path,
                "alt_text": alt_content
            })
        except Exception:
            continue

    if alt_text_map:
        map_file_path = os.path.join(document_dir, "alt_image_map.json")
        with open(map_file_path, "w", encoding="utf-8") as f:
            json.dump(alt_text_map, f, indent=4)

    return images

def extract_images_from_pptx(document_dir, file_path, images_dir):
    prs = Presentation(file_path)
    os.makedirs(images_dir, exist_ok=True)
    images = []
    alt_text_map = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    nvPicPr = shape._element.find(qn('p:nvPicPr'))
                    if nvPicPr is not None:
                        cNvPr = nvPicPr.find(qn('p:cNvPr'))
                        if cNvPr is not None:
                            alt_text = cNvPr.get('descr')
                            if not alt_text or ':' not in alt_text:
                                continue

                            image_name, _, alt_content = alt_text.partition(':')
                            image_name = image_name.strip()
                            alt_content = alt_content.strip()

                            if not image_name or not alt_content:
                                continue
                            image_name = image_name.replace(" ","_")
                            image = shape.image
                            ext = image.ext
                            image_bytes = image.blob
                            img_filename = f"{image_name}.{ext}"
                            img_path = os.path.join(images_dir, img_filename)

                            with open(img_path, "wb") as f:
                                f.write(image_bytes)

                            images.append(img_path)
                            alt_text_map.append({
                                "path": img_path,
                                "alt_text": alt_content
                            })
                except Exception:
                    continue

    if alt_text_map:
        map_file_path = os.path.join(document_dir, "alt_image_map.json")
        with open(map_file_path, "w", encoding="utf-8") as f:
            json.dump(alt_text_map, f, indent=4)

    return images

#===RAG Helper Methods===
# Initialize once globally
def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")

def generate_gpt4_description(image_path):
    image_data = encode_image(image_path)
    response = client.chat.completions.create(
        model="gpt-4-turbo",
        messages=[
            {"role": "system", "content": "Describe the image in great detail, including objects, people, actions, and background elements."},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Describe this image in full detail."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_data}"}}
                ],
            },
        ],
    )
    return response.choices[0].message.content

def chunk_text(text, chunk_size=500, overlap=250):
    tokens = text.split()
    if chunk_size <= overlap:
        overlap = chunk_size // 2
    step_size = max(1, chunk_size - overlap)
    chunks = [" ".join(tokens[i:i+chunk_size]) for i in range(0, len(tokens), step_size)]
    return chunks